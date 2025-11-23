try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: 'python-pptx' library is not installed.")
    print("Please install it using: pip install python-pptx")
    exit(1)

def create_presentation():
    prs = Presentation()

    # Helper to add a slide with title and content
    def add_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1] # Bulleted List Layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        
        # Content
        tf = slide.placeholders[1].text_frame
        tf.text = content_text_list[0] 
        
        for item in content_text_list[1:]:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

    # Helper for Title Slide
    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0] # Title Slide Layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title_text
        subtitle.text = subtitle_text

    # --- Slide 1: Title ---
    add_title_slide("GameDay-Relics", "A Modern E-Commerce Platform for Sports Memorabilia")

    # --- Slide 2: Introduction ---
    add_slide("Introduction", [
        "What is GameDay-Relics?",
        "A full-stack e-commerce application designed for buying and selling sports relics and memorabilia.",
        "Connects buyers with sellers in a secure and user-friendly environment.",
        "Built with modern web technologies for performance and scalability."
    ])

    # --- Slide 3: Technology Stack ---
    # Custom layout for this one to split Frontend/Backend if possible, but keeping simple for now
    add_slide("Technology Stack", [
        "Frontend:",
        "- Framework: React (via Vite) for a fast and reactive UI.",
        "- Language: TypeScript for type safety and code quality.",
        "- Styling: Tailwind CSS for modern, responsive design.",
        "- Routing: React Router DOM for seamless navigation.",
        "",
        "Backend:",
        "- Runtime: Node.js.",
        "- Framework: Express.js for robust API handling.",
        "- Database: MongoDB (with Mongoose) for flexible data modeling."
    ])

    # --- Slide 4: Key Integrations & Tools ---
    add_slide("Key Integrations & Tools", [
        "Security & Payments:",
        "- Authentication: JWT & bcrypt for secure user management.",
        "- Payments: Stripe integration for secure transactions.",
        "",
        "Media & Utilities:",
        "- File Storage: Cloudinary & Multer for efficient image uploads.",
        "- Notifications: React Toastify for real-time user feedback.",
        "- Icons: Lucide React for a clean and consistent icon set."
    ])

    # --- Slide 5: Core Features - User Experience ---
    add_slide("Core Features - User Experience", [
        "Product Browsing: Intuitive catalog with search and filtering.",
        "Shopping Cart: Persistent cart functionality.",
        "Secure Checkout: Integrated Stripe payment gateway.",
        "Responsive Design: Optimized for desktop, tablet, and mobile."
    ])

    # --- Slide 6: Advanced Features - Platform Integrity ---
    add_slide("Advanced Features - Platform Integrity", [
        "Role-Based Access Control: Distinct roles for Admins, Sellers, and Buyers.",
        "Seller Verification: Dedicated flows to verify seller identity.",
        "Dispute Resolution: Built-in system to handle order disputes.",
        "Audit Logging: Comprehensive tracking of system activities."
    ])

    # --- Slide 7: Backend Architecture ---
    add_slide("Backend Architecture", [
        "RESTful API: Organized routes for Users, Products, Orders, etc.",
        "Data Models:",
        "- Users (Buyers/Sellers/Admins)",
        "- Products (Inventory)",
        "- Orders (Transactions)",
        "- Disputes & Verifications",
        "Middleware: Custom auth, error handling, and file processing."
    ])

    # --- Slide 8: Admin Capabilities ---
    add_slide("Admin Capabilities", [
        "User Management: Oversee user accounts and roles.",
        "Product Oversight: Monitor listings and ensure quality.",
        "Dispute Management: Intervene in transaction issues.",
        "System Monitoring: Access to audit logs and system health."
    ])

    # --- Slide 9: Conclusion ---
    add_slide("Conclusion", [
        "GameDay-Relics delivers:",
        "- A secure, scalable, and feature-rich marketplace.",
        "- A modern developer experience with TypeScript and Vite.",
        "- Trust and safety features like verification and dispute handling.",
        "",
        "Thank You!"
    ])

    output_file = "GameDay-Relics-Presentation.pptx"
    prs.save(output_file)
    print(f"Successfully created presentation: {output_file}")

if __name__ == "__main__":
    create_presentation()
